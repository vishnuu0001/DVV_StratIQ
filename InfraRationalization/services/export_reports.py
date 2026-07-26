# ---------------------------------------------------------------------------
# Author: Vishnuu A
# Scope: services/export_reports.py
# Date: 2025-09-23
# ---------------------------------------------------------------------------
"""
services/export_reports.py
Export & Executive Reporting.

Provides:
  - PDF report: executive summary, risk heatmap, key recommendations (reportlab)
  - Excel export: full asset register with all rationalization fields (openpyxl)
  - PowerPoint export: summary slide deck (python-pptx)
  - CMDB CSV export: ServiceNow-compatible format
"""
from __future__ import annotations

import csv
import io
import logging
from datetime import datetime
from typing import Any

log = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════════════════════════
# EXCEL EXPORT
# ═══════════════════════════════════════════════════════════════════════════════

# Function: export_excel
def export_excel(report: dict) -> bytes:
    """Generate full asset register as .xlsx bytes."""
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
    except ImportError:
        raise RuntimeError("openpyxl not installed. Run: pip install openpyxl")

    wb = openpyxl.Workbook()

    # ── Sheet 1: Summary ─────────────────────────────────────────────────────
    ws_summary = wb.active
    ws_summary.title = "Executive Summary"
    _write_summary_sheet(ws_summary, report)

    # ── Sheet 2: Server Asset Register ───────────────────────────────────────
    ws_servers = wb.create_sheet("Asset Register")
    _write_server_sheet(ws_servers, report)

    # ── Sheet 3: Software Inventory ──────────────────────────────────────────
    ws_sw = wb.create_sheet("Software Inventory")
    _write_software_sheet(ws_sw, report)

    # ── Sheet 4: Storage Details ─────────────────────────────────────────────
    ws_storage = wb.create_sheet("Storage Details")
    _write_storage_sheet(ws_storage, report)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# Function: _xl_header_style
def _xl_header_style(ws, row: int, col_count: int):
    """Apply dark header style to a row."""
    try:
        from openpyxl.styles import Font, PatternFill, Alignment
        fill = PatternFill("solid", fgColor="1E293B")
        font = Font(color="FFFFFF", bold=True, size=10)
        align = Alignment(horizontal="center", vertical="center", wrap_text=True)
        for col in range(1, col_count + 1):
            cell = ws.cell(row=row, column=col)
            cell.fill = fill
            cell.font = font
            cell.alignment = align
    except Exception:
        pass


# Function: _write_summary_sheet
def _write_summary_sheet(ws, report: dict):
    sections = report.get("sections") or {}
    ca = sections.get("cloud_assessment") or {}
    cr = sections.get("cloud_readiness") or {}
    gen_at = report.get("generated_at") or datetime.utcnow().isoformat()

    ws.column_dimensions["A"].width = 35
    ws.column_dimensions["B"].width = 30

    rows = [
        ["Infrastructure Rationalization Report", ""],
        ["Report Name",   report.get("report_name", "")],
        ["Generated At",  gen_at],
        ["Provider",      report.get("provider", "")],
        ["Scan ID",       report.get("scan_id", "")],
        [],
        ["== Cloud Assessment ==", ""],
        ["Total Servers",           ca.get("total_servers", 0)],
        ["Total CPU Cores",         ca.get("total_cpu_cores", 0)],
        ["Total RAM (GB)",          ca.get("total_ram_gb", 0)],
        ["Total Storage (TB)",      ca.get("total_storage_tb", 0)],
        [],
        ["== Cloud Readiness ==", ""],
        ["Lift & Shift",            cr.get("lift_and_shift", 0)],
        ["Smart Shift",             cr.get("smart_shift", 0)],
        ["PaaS Shift",              cr.get("paas_shift", 0)],
        ["Decommission",            cr.get("decommission", 0)],
    ]
    for r, row_data in enumerate(rows, start=1):
        for c, val in enumerate(row_data, start=1):
            ws.cell(row=r, column=c, value=val)


_SERVER_COLUMNS = [
    ("Server Name",          lambda s: s.get("server_name") or s.get("name") or ""),
    ("IP Address",           lambda s: s.get("ip_address") or s.get("ip") or s.get("server_ip") or ""),
    ("Hostname",             lambda s: s.get("hostname") or ""),
    ("Environment",          lambda s: s.get("environment") or ""),
    ("Cloud Provider",       lambda s: s.get("cloud_provider") or "onprem"),
    ("Region",               lambda s: s.get("region") or ""),
    ("Resource Group",       lambda s: s.get("resource_group") or ""),
    ("OS Name",              lambda s: s.get("os_name") or s.get("operating_system") or s.get("os") or ""),
    ("OS Family",            lambda s: s.get("os_family") or ""),
    ("OS Version",           lambda s: s.get("os_version") or ""),
    ("OS EOS Date",          lambda s: s.get("os_end_of_support") or ""),
    ("Server Type",          lambda s: s.get("server_type") or ""),
    ("Architecture",         lambda s: s.get("architecture") or s.get("architecture_type") or s.get("compute_hardware_arch") or ""),
    ("CPU Cores",            lambda s: s.get("cpu_cores") or 0),
    ("RAM (GB)",             lambda s: s.get("ram_gb") or s.get("memory_gb") or 0),
    ("Total Storage (GB)",   lambda s: s.get("total_storage_gb") or s.get("internal_storage_gb") or 0),
    ("Storage Type",         lambda s: s.get("storage_type") or ""),
    ("CPU Util (%)",         lambda s: s.get("cpu_util_pct") if (s.get("cpu_util_pct") or -1) >= 0 else ""),
    ("RAM Util (%)",         lambda s: s.get("ram_util_pct") if (s.get("ram_util_pct") or -1) >= 0 else ""),
    ("Utilization Band",     lambda s: s.get("utilization_band") or s.get("utilization") or ""),
    ("Instance Type",        lambda s: s.get("instance_type") or ""),
    ("Boot Type",            lambda s: s.get("boot_type") or ""),
    ("Virtualization State", lambda s: s.get("virtualization_state") or ""),
    ("Migration Strategy",   lambda s: s.get("migration_strategy") or ""),
    ("Cloud Ready",          lambda s: str(s.get("cloud_ready", ""))),
    ("Cloud Suitability",    lambda s: s.get("cloud_suitability") or ""),
    ("HA/DR Requirements",   lambda s: s.get("ha_dr_requirements") or ""),
    ("RTO Requirements",     lambda s: s.get("rto_requirements") or ""),
    ("RPO Requirements",     lambda s: s.get("rpo_requirements") or ""),
    ("DB Engine",            lambda s: s.get("db_engine") or ""),
    ("Business Owner",       lambda s: s.get("business_owner") or ""),
    ("Platform Host",        lambda s: s.get("platform_host") or ""),
    ("Workloads",            lambda s: "; ".join(
        f"{w.get('name','')}{' ' + w.get('version','') if w.get('version') else ''}"
        for w in (s.get("workloads") or [])
    )),
    ("License Type",         lambda s: s.get("license_type") or ""),
    ("App Stability",        lambda s: s.get("application_stability") or ""),
    ("Mainframe Dependency", lambda s: s.get("mainframe_dependency") or ""),
    ("Desktop Dependency",   lambda s: s.get("desktop_dependency") or ""),
    ("Deployment Geography", lambda s: s.get("deployment_geography") or ""),
    ("Power (kWh/month)",    lambda s: s.get("power_consumption_kw_month") or ""),
]


# Function: _write_server_sheet
def _write_server_sheet(ws, report: dict):
    servers = report.get("servers") or []
    headers = [col[0] for col in _SERVER_COLUMNS]
    for c, h in enumerate(headers, start=1):
        ws.cell(row=1, column=c, value=h)
    _xl_header_style(ws, 1, len(headers))
    for col_idx in range(1, len(headers) + 1):
        ws.column_dimensions[ws.cell(row=1, column=col_idx).column_letter].width = 22

    for r, srv in enumerate(servers, start=2):
        for c, (_, extractor) in enumerate(_SERVER_COLUMNS, start=1):
            try:
                ws.cell(row=r, column=c, value=extractor(srv))
            except Exception:
                ws.cell(row=r, column=c, value="")


# Function: _write_software_sheet
def _write_software_sheet(ws, report: dict):
    headers = ["Server Name", "Software Name", "Version", "Vendor", "Category",
               "License Type", "Install Date", "Is EOS", "EOS Date", "Days to EOS"]
    for c, h in enumerate(headers, start=1):
        ws.cell(row=1, column=c, value=h)
    _xl_header_style(ws, 1, len(headers))

    row = 2
    for srv in (report.get("servers") or []):
        srv_name = srv.get("server_name") or srv.get("name") or ""
        for sw in (srv.get("installed_software") or []):
            ws.cell(row=row, column=1, value=srv_name)
            ws.cell(row=row, column=2, value=sw.get("name") or "")
            ws.cell(row=row, column=3, value=sw.get("version") or "")
            ws.cell(row=row, column=4, value=sw.get("vendor") or "")
            ws.cell(row=row, column=5, value=sw.get("category") or "")
            ws.cell(row=row, column=6, value=sw.get("license_type") or "")
            ws.cell(row=row, column=7, value=sw.get("install_date") or "")
            ws.cell(row=row, column=8, value=str(sw.get("is_eos", False)))
            ws.cell(row=row, column=9, value=sw.get("eos_date") or "")
            ws.cell(row=row, column=10, value=sw.get("days_to_eos") or "")
            row += 1


# Function: _write_storage_sheet
def _write_storage_sheet(ws, report: dict):
    headers = ["Server Name", "Mount Point", "Size (GB)", "Used (GB)", "Free (GB)", "Disk Type", "IOPS"]
    for c, h in enumerate(headers, start=1):
        ws.cell(row=1, column=c, value=h)
    _xl_header_style(ws, 1, len(headers))

    row = 2
    for srv in (report.get("servers") or []):
        srv_name = srv.get("server_name") or srv.get("name") or ""
        for disk in (srv.get("disks") or []):
            size  = disk.get("size_gb") or 0
            used  = disk.get("used_gb") or 0
            free  = round(size - used, 2) if size > 0 else ""
            ws.cell(row=row, column=1, value=srv_name)
            ws.cell(row=row, column=2, value=disk.get("mount_point") or "")
            ws.cell(row=row, column=3, value=size)
            ws.cell(row=row, column=4, value=used)
            ws.cell(row=row, column=5, value=free)
            ws.cell(row=row, column=6, value=disk.get("disk_type") or "")
            ws.cell(row=row, column=7, value=disk.get("iops") or "")
            row += 1


# ═══════════════════════════════════════════════════════════════════════════════
# CMDB CSV EXPORT (ServiceNow-compatible)
# ═══════════════════════════════════════════════════════════════════════════════

_CMDB_COLUMNS = [
    ("name",                      lambda s: s.get("server_name") or s.get("name") or ""),
    ("ip_address",                lambda s: s.get("ip_address") or s.get("ip") or s.get("server_ip") or ""),
    ("fqdn",                      lambda s: s.get("hostname") or ""),
    ("os",                        lambda s: s.get("os_name") or s.get("operating_system") or ""),
    ("os_version",                lambda s: s.get("os_version") or ""),
    ("cpu_count",                 lambda s: s.get("cpu_cores") or 0),
    ("ram",                       lambda s: str(int((s.get("ram_gb") or s.get("memory_gb") or 0) * 1024)) + " MB"),
    ("disk_space",                lambda s: str(int(s.get("total_storage_gb") or 0)) + " GB"),
    ("virtual",                   lambda s: "true" if (s.get("server_type") or "").lower() == "virtual" else "false"),
    ("environment",               lambda s: s.get("environment") or ""),
    ("cloud_provider",            lambda s: s.get("cloud_provider") or ""),
    ("region",                    lambda s: s.get("region") or ""),
    ("assignment_group",          lambda s: s.get("business_owner") or ""),
    ("u_migration_strategy",      lambda s: s.get("migration_strategy") or ""),
    ("u_cloud_suitability",       lambda s: s.get("cloud_suitability") or ""),
    ("u_ha_dr_requirements",      lambda s: s.get("ha_dr_requirements") or ""),
    ("u_rto_requirements",        lambda s: s.get("rto_requirements") or ""),
    ("u_rpo_requirements",        lambda s: s.get("rpo_requirements") or ""),
    ("u_os_eos_date",             lambda s: s.get("os_end_of_support") or ""),
    ("u_db_engine",               lambda s: s.get("db_engine") or ""),
    ("u_instance_type",           lambda s: s.get("instance_type") or ""),
    ("u_storage_type",            lambda s: s.get("storage_type") or ""),
    ("u_application_stability",   lambda s: s.get("application_stability") or ""),
    ("u_cpu_utilization_pct",     lambda s: str(s.get("cpu_util_pct") or "")),
    ("u_ram_utilization_pct",     lambda s: str(s.get("ram_util_pct") or "")),
    ("u_utilization_band",        lambda s: s.get("utilization_band") or s.get("utilization") or ""),
    ("u_deployment_geography",    lambda s: s.get("deployment_geography") or ""),
    ("u_power_kwh_month",         lambda s: str(s.get("power_consumption_kw_month") or "")),
]


# Function: export_cmdb_csv
def export_cmdb_csv(report: dict) -> bytes:
    """Generate ServiceNow-compatible CMDB CSV as bytes."""
    servers = report.get("servers") or []
    buf = io.StringIO()
    writer = csv.writer(buf, lineterminator="\n")
    writer.writerow([col[0] for col in _CMDB_COLUMNS])
    for srv in servers:
        row = []
        for _, extractor in _CMDB_COLUMNS:
            try:
                row.append(extractor(srv))
            except Exception:
                row.append("")
        writer.writerow(row)
    return buf.getvalue().encode("utf-8")


# ═══════════════════════════════════════════════════════════════════════════════
# PDF EXPORT
# ═══════════════════════════════════════════════════════════════════════════════

# Function: export_pdf
def export_pdf(report: dict) -> bytes:
    """Generate executive summary PDF using reportlab."""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.lib.units import cm
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
        )
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    except ImportError:
        raise RuntimeError("reportlab not installed. Run: pip install reportlab")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    story = []

    # Title
    title_style = ParagraphStyle("Title", parent=styles["Title"],
                                  fontSize=22, textColor=colors.HexColor("#1E293B"),
                                  spaceAfter=6)
    story.append(Paragraph("Infrastructure Rationalization Report", title_style))
    story.append(Paragraph(
        f"Report: {report.get('report_name','')} &nbsp; | &nbsp; Generated: {report.get('generated_at','')[:10]}",
        styles["Normal"]
    ))
    story.append(Spacer(1, 0.5*cm))
    story.append(HRFlowable(width="100%", thickness=2, color=colors.HexColor("#10B981")))
    story.append(Spacer(1, 0.5*cm))

    sections = report.get("sections") or {}
    ca = sections.get("cloud_assessment") or {}
    cr = sections.get("cloud_readiness") or {}
    servers = report.get("servers") or []

    # Executive Summary Stats
    h2 = ParagraphStyle("H2", parent=styles["Heading2"],
                          fontSize=13, textColor=colors.HexColor("#10B981"),
                          spaceBefore=12, spaceAfter=4)
    story.append(Paragraph("Executive Summary", h2))

    summary_data = [
        ["Metric", "Value"],
        ["Total Servers",        str(ca.get("total_servers", len(servers)))],
        ["Total CPU Cores",      str(ca.get("total_cpu_cores", 0))],
        ["Total RAM (GB)",       str(ca.get("total_ram_gb", 0))],
        ["Total Storage (TB)",   str(round(ca.get("total_storage_tb", 0), 2))],
        ["Cloud Ready",          str(cr.get("cloud_ready", 0))],
        ["Lift & Shift",         str(cr.get("lift_and_shift", 0))],
        ["Smart Shift",          str(cr.get("smart_shift", 0))],
        ["Decommission",         str(cr.get("decommission", 0))],
    ]
    tbl = Table(summary_data, colWidths=[8*cm, 8*cm])
    tbl.setStyle(TableStyle([
        ("BACKGROUND",  (0, 0), (-1, 0), colors.HexColor("#1E293B")),
        ("TEXTCOLOR",   (0, 0), (-1, 0), colors.white),
        ("FONTNAME",    (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE",    (0, 0), (-1, -1), 10),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#F8FAFC"), colors.white]),
        ("GRID",        (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
        ("TOPPADDING",  (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING",(0, 0), (-1, -1), 6),
    ]))
    story.append(tbl)
    story.append(Spacer(1, 0.5*cm))

    # Server inventory table (first 20)
    story.append(Paragraph("Server Inventory (Top 20)", h2))
    srv_headers = ["Server Name", "OS", "CPU", "RAM (GB)", "Migration Strategy", "Utilization"]
    srv_rows = [srv_headers]
    for s in servers[:20]:
        srv_rows.append([
            s.get("server_name") or s.get("name") or "",
            (s.get("os_name") or s.get("os") or "")[:30],
            str(s.get("cpu_cores") or 0),
            str(s.get("ram_gb") or s.get("memory_gb") or 0),
            (s.get("migration_strategy") or "").replace("_", " ")[:20],
            s.get("utilization_band") or s.get("utilization") or "",
        ])
    srv_tbl = Table(srv_rows, colWidths=[4.5*cm, 4*cm, 1.5*cm, 2*cm, 3.5*cm, 2.5*cm])
    srv_tbl.setStyle(TableStyle([
        ("BACKGROUND",  (0, 0), (-1, 0), colors.HexColor("#1E293B")),
        ("TEXTCOLOR",   (0, 0), (-1, 0), colors.white),
        ("FONTNAME",    (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE",    (0, 0), (-1, -1), 8),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#F8FAFC"), colors.white]),
        ("GRID",        (0, 0), (-1, -1), 0.3, colors.HexColor("#CBD5E1")),
        ("TOPPADDING",  (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING",(0, 0), (-1, -1), 4),
    ]))
    story.append(srv_tbl)

    # EOS Advisory
    eos_os = sections.get("eos_advisory_os") or []
    if eos_os:
        story.append(Spacer(1, 0.5*cm))
        story.append(Paragraph("OS End-of-Support Advisory", h2))
        warn_style = ParagraphStyle("Warn", parent=styles["Normal"], textColor=colors.HexColor("#EF4444"))
        for eos in eos_os[:10]:
            story.append(Paragraph(
                f"• <b>{eos.get('server_name','')}</b> ({eos.get('server_ip','')}): "
                f"{eos.get('os','')} — EOS: {eos.get('end_of_support','')}",
                warn_style
            ))

    story.append(Spacer(1, 0.5*cm))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#CBD5E1")))
    story.append(Paragraph(
        "Generated by InfraRationalization Module — StratApp Platform",
        styles["Normal"]
    ))

    doc.build(story)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# POWERPOINT EXPORT
# ═══════════════════════════════════════════════════════════════════════════════

# Function: export_pptx
def export_pptx(report: dict) -> bytes:
    """Generate executive PowerPoint slide deck using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK  = RGBColor(0x1E, 0x29, 0x3B)
    GREEN = RGBColor(0x10, 0xB9, 0x81)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    blank_layout = prs.slide_layouts[6]  # blank

    # Function: _add_text
    def _add_text(slide, text, left, top, width, height, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
        from pptx.util import Inches, Pt
        txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf  = txb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txb

    # Function: _add_rect
    def _add_rect(slide, left, top, width, height, color):
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        shape = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # ── Slide 1: Title ───────────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)
    _add_rect(slide1, 0, 0, 13.33, 7.5, DARK)
    _add_rect(slide1, 0, 5.8, 13.33, 0.4, GREEN)
    _add_text(slide1, "Infrastructure Rationalization Report",
              0.5, 2.0, 12, 1.5, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(slide1, report.get("report_name", ""),
              0.5, 3.5, 12, 0.8, size=18, color=GREEN, align=PP_ALIGN.CENTER)
    gen_at = (report.get("generated_at") or "")[:10]
    _add_text(slide1, f"Generated: {gen_at}",
              0.5, 4.4, 12, 0.6, size=14, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Slide 2: Executive Summary Stats ─────────────────────────────────────
    sections = report.get("sections") or {}
    ca = sections.get("cloud_assessment") or {}
    cr = sections.get("cloud_readiness") or {}
    servers = report.get("servers") or []

    slide2 = prs.slides.add_slide(blank_layout)
    _add_rect(slide2, 0, 0, 13.33, 1.0, DARK)
    _add_text(slide2, "Executive Summary", 0.3, 0.15, 12, 0.7,
              size=22, bold=True, color=WHITE)

    stats = [
        ("Total Servers",     str(ca.get("total_servers", len(servers)))),
        ("Total CPU Cores",   str(ca.get("total_cpu_cores", 0))),
        ("Total RAM (GB)",    str(ca.get("total_ram_gb", 0))),
        ("Storage (TB)",      str(round(ca.get("total_storage_tb", 0), 2))),
        ("Cloud Ready",       str(cr.get("cloud_ready", 0))),
        ("Lift & Shift",      str(cr.get("lift_and_shift", 0))),
        ("Smart Shift",       str(cr.get("smart_shift", 0))),
        ("Decommission",      str(cr.get("decommission", 0))),
    ]
    for i, (label, val) in enumerate(stats):
        col = i % 4
        row = i // 4
        x = 0.3 + col * 3.2
        y = 1.3 + row * 2.2
        _add_rect(slide2, x, y, 2.9, 1.8, RGBColor(0x0F, 0x17, 0x2A))
        _add_text(slide2, val, x, y + 0.3, 2.9, 0.8, size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        _add_text(slide2, label, x, y + 1.1, 2.9, 0.5, size=11, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Slide 3: Migration Strategy Breakdown ─────────────────────────────────
    slide3 = prs.slides.add_slide(blank_layout)
    _add_rect(slide3, 0, 0, 13.33, 1.0, DARK)
    _add_text(slide3, "Migration Strategy Breakdown", 0.3, 0.15, 12, 0.7,
              size=22, bold=True, color=WHITE)

    strategy_counts: dict[str, int] = {}
    for s in servers:
        strat = (s.get("migration_strategy") or "Unknown").replace("_", " ").title()
        strategy_counts[strat] = strategy_counts.get(strat, 0) + 1

    y_pos = 1.2
    for strat, count in sorted(strategy_counts.items(), key=lambda x: -x[1]):
        total = len(servers) or 1
        pct   = round(count / total * 100, 1)
        bar_w = max(0.1, pct / 100 * 9)
        _add_text(slide3, f"{strat}", 0.3, y_pos, 3.5, 0.35, size=12, color=DARK)
        _add_rect(slide3, 3.8, y_pos + 0.05, bar_w, 0.25, GREEN)
        _add_text(slide3, f"{count} ({pct}%)", 3.9 + bar_w, y_pos, 3, 0.35, size=11, color=DARK)
        y_pos += 0.5

    # ── Slide 4: Top EOS Risks ────────────────────────────────────────────────
    eos_os = sections.get("eos_advisory_os") or []
    if eos_os:
        slide4 = prs.slides.add_slide(blank_layout)
        _add_rect(slide4, 0, 0, 13.33, 1.0, DARK)
        _add_text(slide4, "OS End-of-Support Risks", 0.3, 0.15, 12, 0.7,
                  size=22, bold=True, color=WHITE)
        y_pos = 1.2
        RED = RGBColor(0xEF, 0x44, 0x44)
        for eos in eos_os[:8]:
            _add_text(slide4,
                      f"• {eos.get('server_name','')} — {eos.get('os','')} — EOS: {eos.get('end_of_support','')}",
                      0.5, y_pos, 12, 0.4, size=12, color=RED)
            y_pos += 0.5

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
